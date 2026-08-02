#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""Cross-check downstream statistics against the result manifest.

把交付文档（论文/报告/PPT）里出现的统计量（CI、P 值）与
results/results.yaml 双向比对；旧项目可读取 07_paper/results.yaml。
  方向A 文中→源：文中每个 CI/P 统计量必须能在 results.yaml 的 display 或旧版 rendered 中找到同值匹配；
                找不到 = 疑似手敲/陈旧/下游私改未回写（高信号）。
  方向B 源→文中：results.yaml 里有、却没在任何交付文档出现的结果（信息，可能是漏用或仅内部）。
匹配前对标点做归一（全角→半角、− → -、去空格、P 小写），只比"值"不比格式，
故同值不同标点不会误报，值不同才报。

用法：
  python check_consistency.py [项目根=.] [--yaml results/results.yaml]
退出码：发现数字不一致 → 1；全过 → 0。
"""
import sys, os, re, glob, argparse

try:
    import yaml
except ImportError:
    print(
        "缺少依赖 pyyaml；请先选择要使用的 Python 环境和安装方式。"
        "本检查器不会自动安装或升级依赖。"
    )
    sys.exit(2)


def norm(s):
    """归一化一个统计量串，只保留可比的值。"""
    if s is None:
        return ""
    s = str(s)
    table = {"：": ":", "，": ",", "（": "(", "）": ")", "＜": "<", "＞": ">",
             "＝": "=", "−": "-", "　": "", "％": "%", "Ｐ": "P", "ｐ": "p"}
    for k, v in table.items():
        s = s.replace(k, v)
    s = s.replace(" ", "").lower()
    return s


# 文中提取：CI 子串 与 P 子串（两类最易出错、最可溯源的统计量）
RE_CI = re.compile(r"\(?95%\s*ci\s*[:：]?\s*[−\-]?\d+\.?\d*\s*[,，]\s*[−\-]?\d+\.?\d*\)?",
                   re.IGNORECASE)
RE_P = re.compile(r"\bp\s*[=＝<＜>＞]\s*0?\.?\d+", re.IGNORECASE)


def extract_text(path):
    ext = path.lower().rsplit(".", 1)[-1]
    try:
        if ext == "md":
            with open(path, encoding="utf-8") as f:
                return f.read()
        if ext == "docx":
            from docx import Document
            d = Document(path)
            parts = [p.text for p in d.paragraphs]
            for t in d.tables:
                for row in t.rows:
                    parts += [c.text for c in row.cells]
            return "\n".join(parts)
        if ext == "pptx":
            from pptx import Presentation
            parts = []
            for s in Presentation(path).slides:
                for sh in s.shapes:
                    if sh.has_text_frame:
                        parts.append(sh.text_frame.text)
                    if sh.has_table:
                        for row in sh.table.rows:
                            parts += [c.text for c in row.cells]
            return "\n".join(parts)
    except Exception as e:
        print(f"  [跳过] 读取失败 {path}: {e}")
    return ""


def source_value_set(doc):
    """Return comparable CI and P values from schema v2 or legacy manifests."""
    ci, pv, full_norms = set(), set(), {}
    for key, r in (doc.get("results") or {}).items():
        rend = r.get("display") or r.get("rendered") or {}
        for which, s in rend.items():
            n = norm(s)
            if not n:
                continue
            full_norms.setdefault(n, key)
            for m in RE_CI.findall(n):
                ci.add(norm(m))
            for m in RE_P.findall(n):
                pv.add(norm(m))
    return ci, pv, full_norms


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("root", nargs="?", default=".")
    ap.add_argument("--yaml", default=None)
    a = ap.parse_args()
    root = a.root
    yaml_path = a.yaml or os.path.join(root, "results", "results.yaml")
    if a.yaml is None and not os.path.exists(yaml_path):
        legacy = os.path.join(root, "07_paper", "results.yaml")
        if os.path.exists(legacy):
            yaml_path = legacy
    if not os.path.exists(yaml_path):
        print(f"找不到结果唯一来源：{yaml_path}"); sys.exit(2)
    with open(yaml_path, encoding="utf-8") as f:
        doc = yaml.safe_load(f) or {}
    src_ci, src_p, full_norms = source_value_set(doc)

    # 待审计交付文档（排除唯一来源、派生 md、备份）
    pats = ["paper/**/*.docx", "paper/**/*.md", "07_paper/**/*.docx", "07_paper/**/*.md",
            "05_reports/**/*.docx", "05_reports/**/*.md", "05_reports/**/*.pptx",
            "04_figures/**/*.pptx", "**/*.pptx"]
    files, seen = [], set()
    for pat in pats:
        for fp in glob.glob(os.path.join(root, pat), recursive=True):
            rp = os.path.normpath(fp)
            base = os.path.basename(rp).lower()
            if rp in seen or "09_backup" in rp.replace("\\", "/"):
                continue
            if base in ("0_result_summaries.md", "results.yaml"):
                continue
            seen.add(rp); files.append(rp)

    problems = 0
    used_keys = set()
    print(f"== 跨文档一致性审计 ==\n源：{yaml_path}（{len(doc.get('results') or {})} 个结果）")
    print(f"待查文档：{len(files)} 个\n")

    for fp in files:
        text = extract_text(fp)
        if not text:
            continue
        ntext = norm(text)
        for n, key in full_norms.items():
            if n and n in ntext:
                used_keys.add(key)
        bad = []
        for m in set(RE_CI.findall(ntext)):
            if norm(m) not in src_ci:
                bad.append(("CI", m))
        for m in set(RE_P.findall(ntext)):
            if norm(m) not in src_p:
                bad.append(("P", m))
        if bad:
            problems += len(bad)
            print(f"[方向A 文中无源匹配] {os.path.relpath(fp, root)}")
            for kind, tok in bad:
                print(f"    {kind}: {tok}  ← 疑似手敲/陈旧/未回写源")

    # 方向B：源里有、文档全未用
    all_keys = set((doc.get("results") or {}).keys())
    unused = sorted(all_keys - used_keys)
    if unused:
        print(f"\n[方向B 源有文档未用]（信息，非必错）：{'、'.join(unused)}")

    print(f"\n== 结果：{'发现 %d 处需处理' % problems if problems else '全部一致，通过'} ==")
    sys.exit(1 if problems else 0)


if __name__ == "__main__":
    main()
