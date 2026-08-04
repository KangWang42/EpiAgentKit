#!/usr/bin/env python3

from __future__ import annotations

import csv
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
DEMO_DIR = ROOT / "docs" / "demo"
RESULT_DIR = DEMO_DIR / "output" / "publication-figures"
OUTPUT_DIR = ROOT / "docs" / "showcase" / "academic-ppt"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

MEETING = {
    "navy": "102A43",
    "teal": "0F766E",
    "teal_light": "E5F3F1",
    "ink": "172033",
    "muted": "607084",
    "line": "D9E2EC",
    "paper": "F7F9FB",
    "white": "FFFFFF",
}

DEFENSE = {
    "wine": "6F263D",
    "wine_dark": "421725",
    "gold": "9C6A2E",
    "gold_light": "E8DCC8",
    "ink": "2D2527",
    "muted": "5F5558",
    "line": "D9CCC8",
    "paper": "F7F2EA",
    "white": "FFFDFC",
}


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def set_background(slide, color: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color)


def add_box(
    slide,
    x,
    y,
    w,
    h,
    *,
    fill: str,
    line: str | None = None,
    radius: bool = False,
):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(0.8)
    if radius:
        shape.adjustments[0] = 0.08
    return shape


def add_text(
    slide,
    text: str,
    x,
    y,
    w,
    h,
    *,
    size: float,
    color: str,
    font: str,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin: float = 0.0,
    line_spacing: float = 1.08,
):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.word_wrap = True
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.line_spacing = line_spacing
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_paragraphs(
    slide,
    items: list[tuple[str, str]],
    x,
    y,
    w,
    h,
    *,
    font: str,
    title_color: str,
    body_color: str,
    title_size: float = 16,
    body_size: float = 14,
    gap: float = 10,
):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.clear()
    frame.margin_left = frame.margin_right = Inches(0)
    frame.margin_top = frame.margin_bottom = Inches(0)
    frame.word_wrap = True
    for index, (heading, body) in enumerate(items):
        heading_p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        heading_p.space_before = Pt(0 if index == 0 else gap)
        heading_p.space_after = Pt(3)
        heading_run = heading_p.add_run()
        heading_run.text = heading
        heading_run.font.name = font
        heading_run.font.size = Pt(title_size)
        heading_run.font.bold = True
        heading_run.font.color.rgb = rgb(title_color)
        body_p = frame.add_paragraph()
        body_p.line_spacing = 1.16
        body_p.space_after = Pt(0)
        body_run = body_p.add_run()
        body_run.text = body
        body_run.font.name = font
        body_run.font.size = Pt(body_size)
        body_run.font.color.rgb = rgb(body_color)
    return box


def add_meeting_header(slide, title: str, section: str, page: int) -> None:
    add_box(
        slide,
        Inches(0.62),
        Inches(0.52),
        Inches(0.08),
        Inches(0.76),
        fill=MEETING["teal"],
    )
    add_text(
        slide,
        section,
        Inches(0.86),
        Inches(0.48),
        Inches(2.0),
        Inches(0.28),
        size=9,
        color=MEETING["teal"],
        font="Microsoft YaHei",
        bold=True,
    )
    add_text(
        slide,
        title,
        Inches(0.86),
        Inches(0.75),
        Inches(11.7),
        Inches(0.6),
        size=25,
        color=MEETING["ink"],
        font="Microsoft YaHei",
        bold=True,
    )
    add_text(
        slide,
        str(page),
        Inches(12.25),
        Inches(7.02),
        Inches(0.4),
        Inches(0.22),
        size=8,
        color=MEETING["muted"],
        font="Arial",
        align=PP_ALIGN.RIGHT,
    )


def add_defense_header(slide, title: str, section: str, page: int) -> None:
    add_text(
        slide,
        section,
        Inches(0.72),
        Inches(0.44),
        Inches(2.2),
        Inches(0.26),
        size=9,
        color=DEFENSE["gold"],
        font="SimSun",
        bold=True,
    )
    add_text(
        slide,
        title,
        Inches(0.72),
        Inches(0.78),
        Inches(11.5),
        Inches(0.58),
        size=24,
        color=DEFENSE["wine_dark"],
        font="SimSun",
        bold=True,
    )
    add_box(
        slide,
        Inches(12.55),
        Inches(0.0),
        Inches(0.78),
        Inches(7.5),
        fill=DEFENSE["wine"],
    )
    add_text(
        slide,
        f"{page:02d}",
        Inches(12.70),
        Inches(6.88),
        Inches(0.38),
        Inches(0.24),
        size=9,
        color=DEFENSE["white"],
        font="Georgia",
        align=PP_ALIGN.CENTER,
    )


def add_line(slide, x1, y1, x2, y2, color: str, width: float = 1.5) -> None:
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width)


def add_arrow(slide, x1, y1, x2, y2, color: str, width: float = 1.5) -> None:
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width)
    tail = OxmlElement("a:tailEnd")
    tail.set("type", "triangle")
    tail.set("w", "sm")
    tail.set("len", "sm")
    line.line._get_or_add_ln().append(tail)


def add_table(
    slide,
    rows: list[list[str]],
    x,
    y,
    widths: list[float],
    row_height: float,
    *,
    header_fill: str,
    stripe_fill: str,
    line_color: str,
    font: str,
    font_size: float,
    text_color: str,
):
    table_shape = slide.shapes.add_table(
        len(rows), len(rows[0]), x, y, sum(Inches(value) for value in widths), Inches(row_height * len(rows))
    )
    table = table_shape.table
    for index, width in enumerate(widths):
        table.columns[index].width = Inches(width)
    for row_index, row in enumerate(rows):
        table.rows[row_index].height = Inches(row_height)
        for col_index, value in enumerate(row):
            cell = table.cell(row_index, col_index)
            cell.margin_left = cell.margin_right = Inches(0.10)
            cell.margin_top = cell.margin_bottom = Inches(0.06)
            cell.fill.solid()
            if row_index == 0:
                cell.fill.fore_color.rgb = rgb(header_fill)
            elif row_index % 2 == 0:
                cell.fill.fore_color.rgb = rgb(stripe_fill)
            else:
                cell.fill.fore_color.rgb = rgb("FFFFFF")
            cell.border_left = None
            frame = cell.text_frame
            frame.clear()
            frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            paragraph = frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.LEFT if col_index == 0 else PP_ALIGN.CENTER
            run = paragraph.add_run()
            run.text = value
            run.font.name = font
            run.font.size = Pt(font_size)
            run.font.bold = row_index == 0
            run.font.color.rgb = rgb("FFFFFF" if row_index == 0 else text_color)
            for border_name in ("top", "bottom"):
                border = getattr(cell, f"_tc").get_or_add_tcPr()
                _ = border_name
    add_line(
        slide,
        x,
        y + Inches(row_height * len(rows)),
        x + sum(Inches(value) for value in widths),
        y + Inches(row_height * len(rows)),
        line_color,
        width=1.0,
    )
    return table_shape


def new_presentation() -> Presentation:
    presentation = Presentation()
    presentation.slide_width = SLIDE_W
    presentation.slide_height = SLIDE_H
    return presentation


def load_results() -> tuple[dict[str, str], list[dict[str, str]]]:
    with (RESULT_DIR / "survival-demo-results.csv").open(
        encoding="utf-8-sig", newline=""
    ) as stream:
        summary = next(csv.DictReader(stream))
    with (RESULT_DIR / "cox-forest-results.csv").open(
        encoding="utf-8-sig", newline=""
    ) as stream:
        model_rows = list(csv.DictReader(stream))
    return summary, model_rows


def build_meeting_deck(summary: dict[str, str], model_rows: list[dict[str, str]]) -> Path:
    prs = new_presentation()
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    set_background(slide, MEETING["paper"])
    add_box(slide, Inches(0), Inches(0), Inches(4.55), Inches(7.5), fill=MEETING["navy"])
    add_box(slide, Inches(4.55), Inches(0), Inches(0.16), Inches(7.5), fill=MEETING["teal"])
    add_text(
        slide,
        "组会汇报｜固定模拟数据",
        Inches(0.72),
        Inches(0.78),
        Inches(3.1),
        Inches(0.3),
        size=10,
        color="7FD1C8",
        font="Microsoft YaHei",
        bold=True,
    )
    add_text(
        slide,
        "固定模拟队列的\n生存分析",
        Inches(0.72),
        Inches(1.55),
        Inches(3.15),
        Inches(1.65),
        size=31,
        color=MEETING["white"],
        font="Microsoft YaHei",
        bold=True,
        line_spacing=1.12,
    )
    add_text(
        slide,
        "从模型口径到可复核图表",
        Inches(0.72),
        Inches(3.45),
        Inches(3.2),
        Inches(0.4),
        size=16,
        color="D9E8F2",
        font="Microsoft YaHei",
    )
    add_text(
        slide,
        "模拟结果不构成医学证据",
        Inches(0.72),
        Inches(6.62),
        Inches(3.1),
        Inches(0.3),
        size=11,
        color="D9E8F2",
        font="Microsoft YaHei",
    )
    for index, (number, label, value) in enumerate(
        (
            ("01", "研究对象", "1,200 人固定模拟队列"),
            ("02", "随访窗口", "36 个月行政随访"),
            ("03", "主要分析", "多变量 Cox 回归"),
        )
    ):
        y = 1.35 + index * 1.55
        add_text(
            slide,
            number,
            Inches(5.35),
            Inches(y),
            Inches(0.65),
            Inches(0.42),
            size=17,
            color=MEETING["teal"],
            font="Georgia",
            bold=True,
        )
        add_text(
            slide,
            label,
            Inches(6.15),
            Inches(y - 0.02),
            Inches(1.5),
            Inches(0.35),
            size=12,
            color=MEETING["muted"],
            font="Microsoft YaHei",
            bold=True,
        )
        add_text(
            slide,
            value,
            Inches(6.15),
            Inches(y + 0.37),
            Inches(5.1),
            Inches(0.42),
            size=19,
            color=MEETING["ink"],
            font="Microsoft YaHei",
            bold=True,
        )

    slide = prs.slides.add_slide(blank)
    set_background(slide, MEETING["paper"])
    add_meeting_header(slide, "固定模拟队列支持可复核的时间到事件分析", "研究设计", 1)
    add_text(
        slide,
        "研究对象",
        Inches(0.9),
        Inches(1.78),
        Inches(1.2),
        Inches(0.3),
        size=11,
        color=MEETING["muted"],
        font="Microsoft YaHei",
        bold=True,
    )
    timeline_y = Inches(3.05)
    add_arrow(slide, Inches(1.0), timeline_y, Inches(6.2), timeline_y, MEETING["teal"], 2.2)
    for x, month, risk in ((1.0, "0", "1,200"), (2.72, "12", "907"), (4.44, "24", "714"), (6.15, "36", "570")):
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x - 0.11),
            timeline_y - Inches(0.11),
            Inches(0.22),
            Inches(0.22),
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = rgb(MEETING["white"])
        circle.line.color.rgb = rgb(MEETING["teal"])
        circle.line.width = Pt(2)
        add_text(
            slide,
            f"{month} 月",
            Inches(x - 0.35),
            Inches(3.38),
            Inches(0.7),
            Inches(0.3),
            size=11,
            color=MEETING["ink"],
            font="Microsoft YaHei",
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            f"风险集 {risk}",
            Inches(x - 0.48),
            Inches(3.77),
            Inches(0.96),
            Inches(0.28),
            size=10,
            color=MEETING["ink"],
            font="Microsoft YaHei",
            align=PP_ALIGN.CENTER,
        )
    add_box(
        slide,
        Inches(7.15),
        Inches(1.72),
        Inches(5.18),
        Inches(4.68),
        fill=MEETING["white"],
        line=MEETING["line"],
        radius=True,
    )
    add_paragraphs(
        slide,
        [
            ("分析对象", "固定随机种子生成 1,200 名模拟对象；630 个结局事件。"),
            ("主要模型", "治疗方案、年龄、性别、疾病分期与生物标志物共同进入 Cox 回归。"),
            ("输出合同", "调整后曲线、95% 置信区间、风险集和全部模型项来自同一次实跑。"),
        ],
        Inches(7.55),
        Inches(2.08),
        Inches(4.35),
        Inches(3.8),
        font="Microsoft YaHei",
        title_color=MEETING["teal"],
        body_color=MEETING["ink"],
        title_size=14,
        body_size=13,
        gap=9,
    )

    slide = prs.slides.add_slide(blank)
    set_background(slide, MEETING["paper"])
    add_meeting_header(slide, "强化方案组的调整后无事件生存率持续较高", "主要结果", 2)
    slide.shapes.add_picture(
        str(RESULT_DIR / "adjusted-survival-paper.png"),
        Inches(2.47),
        Inches(1.42),
        width=Inches(8.4),
    )
    add_text(
        slide,
        "图 1　两种治疗方案的调整后无事件生存曲线及 95% 置信区间",
        Inches(2.15),
        Inches(6.72),
        Inches(9.0),
        Inches(0.3),
        size=10,
        color=MEETING["muted"],
        font="Microsoft YaHei",
        align=PP_ALIGN.CENTER,
    )

    slide = prs.slides.add_slide(blank)
    set_background(slide, MEETING["paper"])
    add_meeting_header(slide, "多变量模型支持一致的效应方向", "模型结果", 3)
    table_rows = [["模型项", "风险比（95% 置信区间）", "P 值"]]
    for row in model_rows:
        p_value = float(row["p_value"])
        table_rows.append(
            [row["label"], row["estimate_label"], "<0.001" if p_value < 0.001 else f"{p_value:.3f}"]
        )
    add_table(
        slide,
        table_rows,
        Inches(1.0),
        Inches(1.75),
        [5.5, 3.7, 1.3],
        0.58,
        header_fill=MEETING["navy"],
        stripe_fill="EDF3F7",
        line_color=MEETING["line"],
        font="Microsoft YaHei",
        font_size=12.5,
        text_color=MEETING["ink"],
    )
    add_text(
        slide,
        f"强化方案与较低的无事件风险相关：HR {float(summary['hazard_ratio']):.2f}（{float(summary['ci_lower']):.2f}～{float(summary['ci_upper']):.2f}）。",
        Inches(1.0),
        Inches(6.25),
        Inches(10.5),
        Inches(0.42),
        size=14,
        color=MEETING["teal"],
        font="Microsoft YaHei",
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    slide = prs.slides.add_slide(blank)
    set_background(slide, MEETING["paper"])
    add_meeting_header(slide, "模型解释受模拟设计与比例风险假设限制", "解释边界", 4)
    add_box(
        slide,
        Inches(0.95),
        Inches(1.7),
        Inches(5.05),
        Inches(4.85),
        fill=MEETING["teal_light"],
        radius=True,
    )
    add_text(
        slide,
        "可支持的解释",
        Inches(1.35),
        Inches(2.12),
        Inches(2.3),
        Inches(0.4),
        size=17,
        color=MEETING["teal"],
        font="Microsoft YaHei",
        bold=True,
    )
    add_text(
        slide,
        "在预设模拟机制和已纳入协变量下，\n强化方案组的调整后风险较低，\n曲线方向与模型估计一致。",
        Inches(1.35),
        Inches(2.78),
        Inches(4.15),
        Inches(2.15),
        size=17,
        color=MEETING["ink"],
        font="Microsoft YaHei",
        bold=True,
        line_spacing=1.28,
    )
    add_text(
        slide,
        f"HR {float(summary['hazard_ratio']):.2f}（95% CI {float(summary['ci_lower']):.2f}～{float(summary['ci_upper']):.2f}）",
        Inches(1.35),
        Inches(5.42),
        Inches(4.15),
        Inches(0.5),
        size=17,
        color=MEETING["teal"],
        font="Arial",
        bold=True,
    )
    add_paragraphs(
        slide,
        [
            ("模拟边界", "固定生成机制不能代表真实人群、治疗选择或未测量混杂。"),
            ("模型边界", "解释依赖比例风险假设、函数形式和协变量测量质量。"),
            ("外部边界", "没有外部队列、真实临床结局或实际决策阈值，不能据此形成医学建议。"),
        ],
        Inches(6.75),
        Inches(1.92),
        Inches(5.35),
        Inches(4.35),
        font="Microsoft YaHei",
        title_color=MEETING["navy"],
        body_color=MEETING["ink"],
        title_size=15,
        body_size=13.5,
        gap=11,
    )

    slide = prs.slides.add_slide(blank)
    set_background(slide, MEETING["paper"])
    add_meeting_header(slide, "后续检查聚焦模型假设与外部适用性", "下一步", 5)
    steps = (
        ("01", "模型诊断", "核对 Schoenfeld 残差、函数形式和高影响观测。"),
        ("02", "稳健性分析", "比较替代调整集、分层模型和关键参数敏感性。"),
        ("03", "外部验证", "在授权真实队列中重新估计性能与适用边界。"),
    )
    for index, (number, heading, body) in enumerate(steps):
        y = 1.75 + index * 1.58
        add_text(
            slide,
            number,
            Inches(1.0),
            Inches(y),
            Inches(0.8),
            Inches(0.5),
            size=24,
            color=MEETING["teal"],
            font="Georgia",
            bold=True,
        )
        add_text(
            slide,
            heading,
            Inches(2.0),
            Inches(y),
            Inches(2.0),
            Inches(0.4),
            size=17,
            color=MEETING["navy"],
            font="Microsoft YaHei",
            bold=True,
        )
        add_text(
            slide,
            body,
            Inches(4.15),
            Inches(y),
            Inches(7.2),
            Inches(0.55),
            size=15,
            color=MEETING["ink"],
            font="Microsoft YaHei",
        )
        if index < len(steps) - 1:
            add_line(
                slide,
                Inches(2.0),
                Inches(y + 0.78),
                Inches(11.4),
                Inches(y + 0.78),
                MEETING["line"],
                0.8,
            )

    output = OUTPUT_DIR / "survival-analysis-meeting.pptx"
    prs.save(output)
    return output


def build_defense_deck() -> Path:
    prs = new_presentation()
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    set_background(slide, DEFENSE["paper"])
    add_box(slide, Inches(0), Inches(0), Inches(0.48), Inches(7.5), fill=DEFENSE["wine"])
    add_box(slide, Inches(9.6), Inches(0), Inches(3.73), Inches(7.5), fill=DEFENSE["wine_dark"])
    add_text(
        slide,
        "开题答辩｜方法学模拟研究",
        Inches(1.1),
        Inches(0.9),
        Inches(4.8),
        Inches(0.35),
        size=11,
        color=DEFENSE["gold"],
        font="SimSun",
        bold=True,
    )
    add_text(
        slide,
        "逆概率失访加权的\n稳健性评估",
        Inches(1.1),
        Inches(1.72),
        Inches(7.7),
        Inches(1.65),
        size=31,
        color=DEFENSE["wine_dark"],
        font="SimSun",
        bold=True,
        line_spacing=1.15,
    )
    add_text(
        slide,
        "预设模拟研究方案，不包含实证结果",
        Inches(1.1),
        Inches(3.78),
        Inches(6.3),
        Inches(0.42),
        size=16,
        color=DEFENSE["muted"],
        font="SimSun",
    )
    add_text(
        slide,
        "研究问题\n方法比较\n质量控制\n研究计划",
        Inches(10.25),
        Inches(1.78),
        Inches(2.2),
        Inches(3.2),
        size=19,
        color=DEFENSE["white"],
        font="SimSun",
        bold=True,
        line_spacing=1.65,
    )
    add_text(
        slide,
        "固定模拟研究方案",
        Inches(1.1),
        Inches(6.62),
        Inches(3.4),
        Inches(0.3),
        size=10,
        color=DEFENSE["muted"],
        font="SimSun",
    )

    slide = prs.slides.add_slide(blank)
    set_background(slide, DEFENSE["paper"])
    add_defense_header(slide, "汇报结构", "答辩导航", 1)
    sections = (
        ("一", "研究问题与目标", "问题、目标与估计对象"),
        ("二", "模拟设计与比较方法", "数据生成、失访机制与分析策略"),
        ("三", "评价指标与质量控制", "偏倚、精度、覆盖率与异常处理"),
        ("四", "研究计划与预期产出", "时间安排、复现材料与适用边界"),
    )
    for index, (number, label, detail) in enumerate(sections):
        y = 1.65 + index * 1.15
        add_text(
            slide,
            number,
            Inches(1.05),
            Inches(y),
            Inches(0.6),
            Inches(0.45),
            size=18,
            color=DEFENSE["gold"],
            font="SimSun",
            bold=True,
        )
        add_text(
            slide,
            label,
            Inches(1.9),
            Inches(y),
            Inches(5.0),
            Inches(0.45),
            size=20,
            color=DEFENSE["wine_dark"],
            font="SimSun",
            bold=True,
        )
        add_text(
            slide,
            detail,
            Inches(7.0),
            Inches(y + 0.02),
            Inches(4.65),
            Inches(0.42),
            size=13,
            color=DEFENSE["muted"],
            font="SimSun",
            align=PP_ALIGN.RIGHT,
        )

    slide = prs.slides.add_slide(blank)
    set_background(slide, DEFENSE["paper"])
    add_defense_header(slide, "失访机制可能改变纵向结局估计", "一　研究问题与目标", 2)
    add_paragraphs(
        slide,
        [
            ("研究问题", "当失访同时依赖既往暴露、协变量和结局历史时，\n完整病例分析可能产生系统偏差。"),
            ("目标", "比较完整病例、逆概率失访加权和多重插补在\n不同数据生成条件下的性能。"),
            ("估计目标", "关注总体平均结局差异，并评价点估计、标准误和\n95% 置信区间。"),
        ],
        Inches(0.95),
        Inches(1.72),
        Inches(5.35),
        Inches(4.8),
        font="SimSun",
        title_color=DEFENSE["wine"],
        body_color=DEFENSE["ink"],
        title_size=15,
        body_size=14,
        gap=11,
    )
    boxes = (
        (7.15, 1.88, "既往协变量"),
        (9.75, 1.88, "结局历史"),
        (8.45, 3.45, "失访概率"),
        (8.45, 5.02, "观察到的结局"),
    )
    for x, y, text_value in boxes:
        add_box(
            slide,
            Inches(x),
            Inches(y),
            Inches(2.15),
            Inches(0.72),
            fill=DEFENSE["white"],
            line=DEFENSE["gold_light"],
            radius=True,
        )
        add_text(
            slide,
            text_value,
            Inches(x),
            Inches(y),
            Inches(2.15),
            Inches(0.72),
            size=14,
            color=DEFENSE["wine_dark"],
            font="SimSun",
            bold=True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )
    add_arrow(slide, Inches(8.15), Inches(2.62), Inches(8.95), Inches(3.43), DEFENSE["gold"], 1.6)
    add_arrow(slide, Inches(10.65), Inches(2.62), Inches(9.95), Inches(3.43), DEFENSE["gold"], 1.6)
    add_arrow(slide, Inches(9.52), Inches(4.18), Inches(9.52), Inches(5.0), DEFENSE["wine"], 1.8)

    slide = prs.slides.add_slide(blank)
    set_background(slide, DEFENSE["paper"])
    add_defense_header(slide, "三类分析策略形成可解释的公平比较", "二　模拟设计与比较方法", 3)
    method_rows = [
        ["方法", "核心处理", "主要风险"],
        ["完整病例分析", "仅分析结局完整者", "失访依赖结局历史时可能偏倚"],
        ["逆概率失访加权", "按持续被观察概率的倒数加权", "权重模型失配或极端权重"],
        ["多重插补", "根据观察历史生成缺失结局", "插补模型与相容性假设"],
    ]
    add_table(
        slide,
        method_rows,
        Inches(0.95),
        Inches(1.78),
        [2.25, 4.0, 4.35],
        0.95,
        header_fill=DEFENSE["wine"],
        stripe_fill="F0E7DE",
        line_color=DEFENSE["line"],
        font="SimSun",
        font_size=13,
        text_color=DEFENSE["ink"],
    )
    add_text(
        slide,
        "所有方法使用相同模拟数据、估计目标、重复次数和评价指标。",
        Inches(1.0),
        Inches(5.95),
        Inches(10.3),
        Inches(0.42),
        size=15,
        color=DEFENSE["wine"],
        font="SimSun",
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    slide = prs.slides.add_slide(blank)
    set_background(slide, DEFENSE["paper"])
    add_defense_header(slide, "六个数据生成因素构成预设实验网格", "二　模拟设计与比较方法", 4)
    factors = (
        ("样本量", "500 / 1,000"),
        ("失访比例", "10% / 30%"),
        ("失访机制", "依赖协变量 / 依赖结局历史"),
        ("效应强度", "较小 / 较大"),
        ("模型设定", "正确 / 遗漏非线性项"),
        ("权重截尾", "无 / 1%–99% / 5%–95%"),
    )
    for index, (label, value) in enumerate(factors):
        col = index % 2
        row = index // 2
        x = 0.95 + col * 5.65
        y = 1.65 + row * 1.55
        add_text(
            slide,
            f"{index + 1:02d}",
            Inches(x),
            Inches(y),
            Inches(0.6),
            Inches(0.35),
            size=13,
            color=DEFENSE["gold"],
            font="Georgia",
            bold=True,
        )
        add_text(
            slide,
            label,
            Inches(x + 0.75),
            Inches(y),
            Inches(1.5),
            Inches(0.35),
            size=15,
            color=DEFENSE["wine_dark"],
            font="SimSun",
            bold=True,
        )
        add_text(
            slide,
            value,
            Inches(x + 0.75),
            Inches(y + 0.45),
            Inches(4.15),
            Inches(0.42),
            size=14,
            color=DEFENSE["ink"],
            font="SimSun",
        )
        add_line(
            slide,
            Inches(x + 0.75),
            Inches(y + 1.03),
            Inches(x + 4.95),
            Inches(y + 1.03),
            DEFENSE["line"],
            0.8,
        )

    slide = prs.slides.add_slide(blank)
    set_background(slide, DEFENSE["paper"])
    add_defense_header(slide, "性能评价同时覆盖偏倚、精度和区间覆盖", "三　评价指标与质量控制", 5)
    metrics = (
        ("偏倚", "E(估计值) − 真值", "判断系统偏离"),
        ("均方根误差", "√E[(估计值 − 真值)²]", "综合偏倚与变异"),
        ("95% 覆盖率", "P(区间包含真值)", "检查区间推断"),
    )
    for index, (label, formula, purpose) in enumerate(metrics):
        x = 0.95 + index * 3.75
        add_text(
            slide,
            label,
            Inches(x),
            Inches(1.78),
            Inches(3.15),
            Inches(0.42),
            size=17,
            color=DEFENSE["wine"],
            font="SimSun",
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_box(
            slide,
            Inches(x),
            Inches(2.52),
            Inches(3.15),
            Inches(1.42),
            fill=DEFENSE["white"],
            line=DEFENSE["gold_light"],
            radius=True,
        )
        add_text(
            slide,
            formula,
            Inches(x + 0.15),
            Inches(2.52),
            Inches(2.85),
            Inches(1.42),
            size=16,
            color=DEFENSE["wine_dark"],
            font="Cambria Math",
            bold=True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )
        add_text(
            slide,
            purpose,
            Inches(x),
            Inches(4.38),
            Inches(3.15),
            Inches(0.42),
            size=13,
            color=DEFENSE["muted"],
            font="SimSun",
            align=PP_ALIGN.CENTER,
        )
    add_text(
        slide,
        "每个情景使用相同重复次数和随机数管理；蒙特卡洛误差随主要指标一并报告。",
        Inches(1.2),
        Inches(5.55),
        Inches(9.8),
        Inches(0.6),
        size=15,
        color=DEFENSE["ink"],
        font="SimSun",
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    slide = prs.slides.add_slide(blank)
    set_background(slide, DEFENSE["paper"])
    add_defense_header(slide, "预设质量控制用于隔离极端权重与模型失配", "三　评价指标与质量控制", 6)
    flow = (
        ("生成数据", "保存真值与随机种子"),
        ("拟合模型", "记录收敛和警告"),
        ("检查权重", "分布、极值与有效样本量"),
        ("汇总性能", "偏倚、RMSE 与覆盖率"),
    )
    for index, (label, detail) in enumerate(flow):
        x = 0.85 + index * 2.85
        add_box(
            slide,
            Inches(x),
            Inches(2.2),
            Inches(2.25),
            Inches(1.5),
            fill=DEFENSE["white"],
            line=DEFENSE["gold_light"],
            radius=True,
        )
        add_text(
            slide,
            label,
            Inches(x + 0.16),
            Inches(2.48),
            Inches(1.93),
            Inches(0.36),
            size=15,
            color=DEFENSE["wine"],
            font="SimSun",
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            detail,
            Inches(x + 0.18),
            Inches(3.0),
            Inches(1.89),
            Inches(0.48),
            size=11,
            color=DEFENSE["muted"],
            font="SimSun",
            align=PP_ALIGN.CENTER,
        )
        if index < len(flow) - 1:
            add_arrow(
                slide,
                Inches(x + 2.3),
                Inches(2.95),
                Inches(x + 2.78),
                Inches(2.95),
                DEFENSE["gold"],
                1.7,
            )
    add_text(
        slide,
        "任何收敛失败、空结果或异常权重都进入情景级诊断，不静默删除后继续汇总。",
        Inches(1.25),
        Inches(4.65),
        Inches(9.6),
        Inches(0.62),
        size=16,
        color=DEFENSE["wine_dark"],
        font="SimSun",
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    slide = prs.slides.add_slide(blank)
    set_background(slide, DEFENSE["paper"])
    add_defense_header(slide, "研究计划以可复现模拟和敏感性分析为主线", "四　研究计划与预期产出", 7)
    phases = (
        ("阶段 1", "方案与代码原型", "第 1–2 月"),
        ("阶段 2", "主模拟与诊断", "第 3–5 月"),
        ("阶段 3", "敏感性分析", "第 6–7 月"),
        ("阶段 4", "论文与材料", "第 8–9 月"),
    )
    add_line(slide, Inches(1.1), Inches(3.0), Inches(11.2), Inches(3.0), DEFENSE["gold"], 2.0)
    for index, (phase, label, timing) in enumerate(phases):
        x = 1.1 + index * 3.35
        marker = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x - 0.12),
            Inches(2.88),
            Inches(0.24),
            Inches(0.24),
        )
        marker.fill.solid()
        marker.fill.fore_color.rgb = rgb(DEFENSE["wine"])
        marker.line.fill.background()
        add_text(
            slide,
            phase,
            Inches(x - 0.75),
            Inches(1.78),
            Inches(1.5),
            Inches(0.35),
            size=12,
            color=DEFENSE["gold"],
            font="SimSun",
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            label,
            Inches(x - 1.0),
            Inches(2.2),
            Inches(2.0),
            Inches(0.46),
            size=15,
            color=DEFENSE["wine_dark"],
            font="SimSun",
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            timing,
            Inches(x - 0.8),
            Inches(3.42),
            Inches(1.6),
            Inches(0.36),
            size=12,
            color=DEFENSE["muted"],
            font="SimSun",
            align=PP_ALIGN.CENTER,
        )
    add_text(
        slide,
        "每个阶段均保留可运行代码、参数配置、情景级诊断和汇总结果，下一阶段只读取已验证输出。",
        Inches(1.25),
        Inches(4.65),
        Inches(9.55),
        Inches(0.85),
        size=16,
        color=DEFENSE["ink"],
        font="SimSun",
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    slide = prs.slides.add_slide(blank)
    set_background(slide, DEFENSE["paper"])
    add_defense_header(slide, "预期产出限定为方法比较与适用边界", "四　研究计划与预期产出", 8)
    outputs = (
        ("方法结果", "三类策略在预设失访机制下的偏倚、精度与覆盖率比较。"),
        ("质量结果", "极端权重、模型失配与截尾规则对推断稳定性的影响。"),
        ("复现材料", "模拟代码、参数配置、诊断记录和可追溯结果汇总。"),
    )
    for index, (heading, body) in enumerate(outputs):
        y = 1.72 + index * 1.45
        add_text(
            slide,
            f"0{index + 1}",
            Inches(1.05),
            Inches(y),
            Inches(0.6),
            Inches(0.42),
            size=17,
            color=DEFENSE["gold"],
            font="Georgia",
            bold=True,
        )
        add_text(
            slide,
            heading,
            Inches(1.85),
            Inches(y),
            Inches(2.0),
            Inches(0.42),
            size=17,
            color=DEFENSE["wine"],
            font="SimSun",
            bold=True,
        )
        add_text(
            slide,
            body,
            Inches(4.0),
            Inches(y),
            Inches(7.3),
            Inches(0.56),
            size=15,
            color=DEFENSE["ink"],
            font="SimSun",
        )
    add_box(
        slide,
        Inches(1.0),
        Inches(6.25),
        Inches(10.35),
        Inches(0.58),
        fill=DEFENSE["wine"],
        radius=True,
    )
    add_text(
        slide,
        "开题阶段不预设方法优劣或结果方向；最终结论服从实际模拟输出。",
        Inches(1.2),
        Inches(6.25),
        Inches(9.95),
        Inches(0.58),
        size=14,
        color=DEFENSE["white"],
        font="SimSun",
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )

    output = OUTPUT_DIR / "missing-data-proposal-defense.pptx"
    prs.save(output)
    return output


def main() -> None:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    summary, model_rows = load_results()
    meeting = build_meeting_deck(summary, model_rows)
    defense = build_defense_deck()
    print(f"Generated {meeting.relative_to(ROOT)}")
    print(f"Generated {defense.relative_to(ROOT)}")


if __name__ == "__main__":
    main()
